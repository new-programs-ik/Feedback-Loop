"""
service.py — HTTP analysis worker around the engine (stateless: no DB writes).

The Next.js app calls this server→server to analyze one class. Given a Vimeo link OR raw
transcript text plus the class context, it fetches the transcript (if needed) and runs the
engine, returning the structured result. Persistence lives in Supabase (owned by the web app).

Run locally:   uvicorn service:app --port 8000
Deploy (free): Render / Cloud Run. Set ANTHROPIC_API_KEY (+ VIMEO_ACCESS_TOKEN) in the env.

Endpoints:
  GET  /health           -> liveness + which capabilities are configured
  POST /dry-run          -> {cues, windows, est_tokens}         (transcript text; no API key)
  POST /transcript       -> {text, video_id, language, chars}   (fetch captions from a Vimeo URL)
  POST /analyze          -> {result, meta, transcript_source}   (needs ANTHROPIC_API_KEY)

Optional shared secret: if WORKER_API_KEY is set, callers must send `Authorization: Bearer <it>`.
"""
from __future__ import annotations

import base64
import io
import json as _json
import logging
import os
from typing import Literal, Optional

from fastapi import BackgroundTasks, Depends, FastAPI, Header, HTTPException
from pydantic import BaseModel, model_validator

import config

config.load_env()

import engine as E  # noqa: E402  (after load_env so config is present)
import store as ST  # noqa: E402
import vimeo as V  # noqa: E402

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s", datefmt="%H:%M:%S")
app = FastAPI(title="Ratings Analysis Worker", version="2.0")

WORKER_API_KEY = os.environ.get("WORKER_API_KEY") or None


def require_worker_auth(authorization: Optional[str] = Header(default=None)) -> None:
    """If WORKER_API_KEY is configured, require it as a Bearer token (server→server)."""
    if not WORKER_API_KEY:
        return
    if authorization != f"Bearer {WORKER_API_KEY}":
        raise HTTPException(status_code=401, detail="invalid worker credentials")


class MaterialFile(BaseModel):
    filename: str = "materials.txt"
    b64: str                                     # base64-encoded file bytes


class AnalyzeRequest(BaseModel):
    transcript: Optional[str] = None
    vimeo_url: Optional[str] = None
    course: str = "(unspecified)"
    cohort: str = "(unspecified)"
    topic: str = "(unspecified)"
    instructor: str = "(unspecified)"
    rating: str = "(unspecified)"
    num_ratings: Optional[int] = None
    agenda: str = "(not provided)"
    class_type: Literal["live_class", "ars"] = "live_class"
    materials_text: str = ""                     # pasted class-materials text (optional)
    materials_files: list[MaterialFile] = []     # one or more uploaded files (slides/notebook/doc)
    # NOTE: materials are used in-memory for this analysis only and are NEVER persisted.

    @model_validator(mode="after")
    def _need_a_source(self):
        has_tx = bool(self.transcript and self.transcript.strip())
        has_url = bool(self.vimeo_url and self.vimeo_url.strip())
        if not has_tx and not has_url:
            raise ValueError("provide either 'transcript' text or a 'vimeo_url'")
        return self

    def context(self) -> str:
        base = E.build_context(self.course, self.topic, self.instructor, self.rating, self.agenda)
        label = ("Assignment Review Session (ARS) — solutions to assigned problems are reviewed and doubts cleared"
                 if self.class_type == "ars" else "Live class (weekly teaching session)")
        return base + f"\nSession type: {label}"


class AnalyzeAsyncRequest(AnalyzeRequest):
    class_id: str                                # the queued class row to write the result back to


class TranscriptRequest(BaseModel):
    vimeo_url: str


# ─────────────────────────── class-materials text extraction ───────────────────────────
def extract_text(filename: str, data: bytes) -> str:
    """Pull plain text out of an uploaded materials file (slides / notebook / doc)."""
    name = (filename or "").lower()
    try:
        if name.endswith((".txt", ".md", ".vtt", ".srt")):
            return data.decode("utf-8", "ignore")
        if name.endswith(".ipynb"):
            nb = _json.loads(data.decode("utf-8", "ignore"))
            parts = []
            for cell in nb.get("cells", []):
                src = "".join(cell.get("source", []))
                parts.append(f"```\n{src}\n```" if cell.get("cell_type") == "code" else src)
            return "\n\n".join(p for p in parts if p.strip())
        if name.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            return "\n\n".join((page.extract_text() or "") for page in reader.pages)
        if name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text_frame.text for sh in slide.shapes
                         if sh.has_text_frame and sh.text_frame.text.strip()]
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
        if name.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except HTTPException:
        raise
    except Exception as e:  # corrupt file, parser error
        raise HTTPException(status_code=422, detail=f"could not read materials file '{filename}': {e}")
    raise HTTPException(status_code=422,
                        detail=f"unsupported materials file type: '{filename}' "
                               "(use .pdf, .pptx, .docx, .txt, .md or .ipynb)")


def gather_materials(req: AnalyzeRequest) -> str:
    """Combine pasted materials text + one or more uploaded materials files into one string.
    Held in memory for this request only — never persisted anywhere."""
    parts = []
    if req.materials_text and req.materials_text.strip():
        parts.append(req.materials_text.strip())
    for mf in req.materials_files:
        try:
            data = base64.b64decode(mf.b64)
        except Exception:
            raise HTTPException(status_code=422, detail=f"materials file '{mf.filename}' is not valid base64")
        text = extract_text(mf.filename or "materials.txt", data)
        if text.strip():
            parts.append(f"=== {mf.filename} ===\n{text}")
    return "\n\n".join(parts)


class ReviseRequest(BaseModel):
    feedback: str                 # the current draft/edited text
    instruction: str              # the PM's plain-English change request
    context: str = ""             # optional class context for grounding
    flags_json: str = ""          # optional verified flags (grounding)


@app.get("/health")
def health() -> dict:
    return {
        "status": "ok",
        "model": E.CFG.model,
        "anthropic_key": bool(os.environ.get("ANTHROPIC_API_KEY")),
        "vimeo_token": bool(os.environ.get("VIMEO_ACCESS_TOKEN")),
    }


@app.post("/dry-run")
def dry_run(req: AnalyzeRequest) -> dict:
    if not (req.transcript and req.transcript.strip()):
        raise HTTPException(status_code=422, detail="dry-run needs 'transcript' text")
    cues = E.parse_cues(req.transcript)
    if not cues:
        raise HTTPException(status_code=422, detail="no cues parsed from transcript")
    return {"cues": len(cues), "windows": len(E.chunk_by_time(cues)), "est_tokens": E.est_tokens(cues)}


@app.post("/transcript", dependencies=[Depends(require_worker_auth)])
def transcript(req: TranscriptRequest) -> dict:
    """Fetch captions from a Vimeo URL only (lets the UI preview before spending on analysis)."""
    try:
        info = V.fetch_transcript(req.vimeo_url)
        return {
            "video_id": info["video_id"],
            "language": info["language"],
            "format": info["format"],
            "chars": len(info["text"]),
            "text": info["text"],
        }
    except V.VimeoNoCaptions as e:
        raise HTTPException(status_code=422, detail=str(e))
    except V.VimeoAuthError as e:
        raise HTTPException(status_code=502, detail=str(e))
    except V.VimeoError as e:
        raise HTTPException(status_code=502, detail=str(e))


def _run_analysis_job(req: AnalyzeAsyncRequest) -> None:
    """The background job: fetch transcript + digest materials + analyze + save to the DB."""
    try:
        transcript_text = req.transcript
        source = "upload"
        if not (transcript_text and transcript_text.strip()):
            info = V.fetch_transcript(req.vimeo_url)  # type: ignore[arg-type]
            transcript_text = info["text"]
            source = "vimeo"
        materials = gather_materials(req)
        result, meta = E.analyse_text(transcript_text, req.context(), req.class_type, materials)
        ST.persist_analysis(req.class_id, result, meta, transcript_text, source)
        logging.info("async analysis stored for class %s (cost $%s)", req.class_id, meta.get("cost_usd"))
    except Exception as e:  # noqa: BLE001 — background job, record failure so the UI can show it
        logging.exception("async analysis failed for class %s", req.class_id)
        ST.mark_failed(req.class_id, str(e))


@app.post("/analyze-async", dependencies=[Depends(require_worker_auth)])
def analyze_async(req: AnalyzeAsyncRequest, background: BackgroundTasks) -> dict:
    """Start the analysis in the background and return immediately (so the web request never times
    out). The worker writes the result straight to the DB when done; the UI polls for it."""
    if not os.environ.get("DATABASE_URL"):
        raise HTTPException(status_code=500, detail="worker has no DATABASE_URL configured for async persistence")
    background.add_task(_run_analysis_job, req)
    return {"status": "accepted", "class_id": req.class_id}


@app.post("/revise", dependencies=[Depends(require_worker_auth)])
def revise(req: ReviseRequest) -> dict:
    """Rewrite a feedback draft per the PM's instruction (the review-page 'fix it' agent)."""
    try:
        text, meta = E.revise_feedback(req.feedback, req.instruction, req.context, req.flags_json)
        return {"feedback": text, "meta": meta}
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        logging.exception("revise failed")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/analyze", dependencies=[Depends(require_worker_auth)])
def analyze(req: AnalyzeRequest) -> dict:
    transcript_text = req.transcript
    source = "upload"
    try:
        if not (transcript_text and transcript_text.strip()):
            info = V.fetch_transcript(req.vimeo_url)  # type: ignore[arg-type]
            transcript_text = info["text"]
            source = "vimeo"
        materials = gather_materials(req)
        result, meta = E.analyse_text(transcript_text, req.context(), req.class_type, materials)
        return {
            "result": result,
            "meta": meta,
            "transcript_source": source,
            "transcript_chars": len(transcript_text),
            "transcript_used": transcript_text,
            "materials_chars": len(materials),
        }
    except HTTPException:
        raise  # e.g. a 422 from materials extraction — don't relabel as 500
    except V.VimeoNoCaptions as e:
        raise HTTPException(status_code=422, detail=str(e))
    except V.VimeoAuthError as e:
        raise HTTPException(status_code=502, detail=str(e))
    except V.VimeoError as e:
        raise HTTPException(status_code=502, detail=str(e))
    except ValueError as e:  # empty/garbled transcript, invalid model output
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:  # auth, network, etc.
        logging.exception("analyze failed")
        raise HTTPException(status_code=500, detail=str(e))
